#This script goes through every sub folder from an identified folder
#It searches for a string within popular document types: office, pdf, archives

import os,re
from PyPDF2 import PdfFileReader
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
import argparse
from pyunpack import Archive
import shutil
import time
startTime = time.time()

#Setting up arguments for command line
parser = argparse.ArgumentParser()
parser.add_argument('-path', '--source_path', help="Takes the directory path in quotes as input")
parser.add_argument('-text', '--text_pattern', help="Takes the text pattern as input")

args = parser.parse_args()

#Assigning variables and adding regular expression for the text
directory = args.source_path
text = args.text_pattern+".+}" #Has the text with one or more characters and ends with }

#Going through each file in the directory to unpack
for root, dirs, files in os.walk(directory):
    for file in files:
        if file.endswith(".zip"):
            out_zip=os.path.join(root, "extracted_zip")
            Archive(file).extractall(out_zip,auto_create_dir=True)
            break

        if file.endswith(".rar"):
            out_rar=os.path.join(root, "extracted_rar")
            Archive(file).extractall(out_rar,auto_create_dir=True)
            break

#Going through each files (unpacked files are included) and finding the text
for root, dirs, files in os.walk(directory):
    for file in files:
        
        if file.endswith(".txt"):
            file_path=os.path.join(root, file)
            data=''
            with open(file_path, 'r') as f:
                data+=f.read()
                result=re.findall(text,data)
                if not result:
                    continue
                else:
                    print(result[0])
                    break    
            
        if file.endswith(".pdf"):
            file_path=os.path.join(root, file)
            with open(file_path, 'rb') as f:
                pdf_file=PdfFileReader(f)
                data=''
                for i in range(pdf_file.numPages):
                    page = pdf_file.getPage(i)
                    data += page.extractText()
                result=re.findall(text,data)
                if not result:
                    continue
                else:
                    print(result[0])
                    break 
                
        if file.endswith(".xlsx"):
            file_path=os.path.join(root, file)
            data=''
            wb = load_workbook(file_path)
            ws = wb.active
            for row in ws:
                for cell in row:
                    if cell.value is not None:
                        data+=str(cell.value)
            result=re.findall(text,data)
            if not result:
                    continue
            else:
                print(result[0])
                break         

        if file.endswith(".pptx"):
            file_path=os.path.join(root, file)
            data=''
            with open(file_path, 'rb') as f:
                prs = Presentation(f)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            data=shape.text
                result=re.findall(text,data)
                if not result:
                    continue
                else:
                    print(result[0])
                    break 
        
        if file.endswith(".docx"):
            file_path=os.path.join(root, file)
            data=''
            with open(file_path, 'rb') as f:
                doc = Document(file)
                data=''
                for para in doc.paragraphs:
                    data += '\n'+ para.text
                result=re.findall(text,data)
                if not result:
                    continue
                else:
                    print(result[0])
                    break 

#Removing unpacked folders

if 'out_zip' in locals() or 'out_rar' in locals():
    ans = input('Do you want to remove all the files that were extracted from compressed directories?   y/n?')
    if ans.lower() =='y':
        #shutil.rmtree(out_zip)
        shutil.rmtree(out_rar)

executionTime = (time.time() - startTime)
print('Execution time in seconds: ' + str(executionTime))