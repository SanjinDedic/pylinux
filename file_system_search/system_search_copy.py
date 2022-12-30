#This script goes through every sub folder from an identified folder
#It searches for a string within popular document types: office, pdf, archives

import os,re
import fitz
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
import argparse
from pyunpack import Archive
import time

startTime = time.time()
counter=0
totalsize=0
counter_txt=0
counter_pdf=0
counter_docx=0
counter_xlsx=0
counter_pptx=0
final=list()


#Setting up arguments for command line
parser = argparse.ArgumentParser()
parser.add_argument('-path', '--source_path', help="Takes the directory path in quotes as input")
parser.add_argument('-text', '--text_pattern', help="Takes the text pattern as input")

args = parser.parse_args()

#Assigning variables and adding regular expression for the text
directory = args.source_path
text = args.text_pattern+".+}" #Has the text with one or more characters and ends with }


def logsize(root,file):
    global totalsize
    size=os.path.getsize(os.path.join(root, file))
    totalsize+=size
    print("{f} | {size} | Bytes".format(f=file,size=size))
    

def progress(percent=0, width=30):
    left = width * percent // 1000
    right = width - left
    
    tags = "#" * left
    spaces = " " * right
    percents = f"{percent:.0f} files done"
    
    print("\r[", tags, spaces, "]", percents, sep="", end="", flush=True)
# Example run

def extract(root,file):
    out_zip=os.path.join(root, "extracted")
    file_path=os.path.join(root, file)
    Archive(file_path).extractall(out_zip,auto_create_dir=True)

def check(root,file):
    global counter,counter_txt,counter_pdf,counter_docx,counter_xlsx,counter_pptx
    
   
    if file.endswith(".txt"):
        file_path=os.path.join(root, file)
        data=''
        
        with open(file_path, 'r') as f:
            data+=f.read()
            result=re.findall(text,data)
            if  result:
                final.append(result[0]+"in file"+file)
        counter_txt+=1

                    
                
    if file.endswith(".pdf"):
        file_path=os.path.join(root, file)
        doc = fitz.open(file_path)
        data=''
        for page in doc:
            data+=page.get_text()
        result=re.findall("data",data)
        if result:
            final.append(result[0]+" in file "+file)
        counter_pdf+=1
            

    if file.endswith(".xlsx"):
        file_path=os.path.join(root, file)
        data=''
        with open(file_path, 'rb') as f:
            wb = load_workbook(f, read_only=True, data_only=True)
        ws = wb.active
        
        values=[str(cell.value) for row in ws for cell in row if cell.value is not None]
        data=' '.join(values)            
        result=re.findall(text,data)
        if result:
            final.append(result[0]+" in file "+file)
        counter_xlsx+=1
            

    if file.endswith(".pptx"):
        file_path=os.path.join(root, file)
        data=''
        prs = Presentation(file_path)
        values=[shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
        data=' '.join(values)
        result=re.findall(text,data)
        if result:
            final.append(result[0]+" in file "+file)
        counter_pptx+=1

    if file.endswith(".docx"):
        file_path=os.path.join(root, file)
        data=''
        doc = Document(file_path)
        data=''
        for para in doc.paragraphs:
            data += '\n'+ para.text
        result=re.findall(text,data)
        if result:
            final.append(result[0]+" in file "+file)
        counter_docx+=1
    counter+=1
    progress(counter)
    

#Going through each file in the directory to unpack
#[extract(root,file) for root, dirs, files in os.walk(directory) for file in files if file.endswith((".zip",".rar")) ]#map(str.upper,os.walk(directory))
#print(all_files)

[check(root,file)  for root, dirs, files in os.walk(directory) for file in files if file.endswith((".txt",".pdf",".xlsx",".pptx",".docx"))]
print("\n")
[print(r) for r in final]
"""
if 'out_zip' in locals() :
    ans = input('Do you want to remove all the files that were extracted from compressed directories?   y/n?')
    if ans.lower() =='y':
        shutil.rmtree(out_zip)
        shutil.rmtree(out_rar)
"""
executionTime = (time.time() - startTime)
print('Execution time in seconds: ' + str(executionTime))
#[logsize(root,file) for root, dirs, files in os.walk(directory) for file in files if file.endswith((".txt",".pdf",".xlsx",".pptx",".docx")) ]
#print("Total text files {txtf}\n Total pdf files {pdff}\n Total docx files {docxf}\n Total xlsx files {xlsxf}\n Total pptx files {pptxf}".format(txtf=counter_txt,pdff=counter_pdf,docxf=counter_docx,xlsxf=counter_xlsx,pptxf=counter_pptx))
#x=[len(dirs) for root, dirs, files in os.walk(directory)]
#print("Total directories {size}".format(size=x[0]))
