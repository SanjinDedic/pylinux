#This script goes through every sub folder from an identified folder
#It searches for a string within popular document types: office, pdf, archives

import os,re
import fitz
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
import argparse
from pyunpack import Archive
import time
import datetime
import shutil

startTime = time.time()
counter,totalsize,counter_txt,counter_pdf,counter_docx,counter_xlsx,counter_pptx,counter_otherfiles=0,0,0,0,0,0,0,0
final=list()
out_zip=''


#Setting up arguments for command line
parser = argparse.ArgumentParser(prog="system_search_noload.py")
parser.add_argument('-path', '--source_path', help="Takes the directory path in quotes as input")
parser.add_argument('-text', '--text_pattern', help="Takes the text pattern as input")
parser.add_argument('-save', '--save_log',action="store_true", help="Saves a log file in a folder")

args = parser.parse_args()

#Assigning variables and adding regular expression for the text
directory = args.source_path
text = args.text_pattern+".+" #Has the text with one or more characters and ends with }


def logsize(root,file):
    global totalsize
    size=os.path.getsize(os.path.join(root, file))
    totalsize+=size
    return "{f} | {size} | Bytes".format(f=file,size=size)
    
def other(file):
    global counter_otherfiles
    counter_otherfiles+=1


def extract(root,file):
    global out_zip
    out_zip=os.path.join(root, "extracted")
    file_path=os.path.join(root, file)
    Archive(file_path).extractall(out_zip,auto_create_dir=True)

def check(root,file):
    global counter,counter_txt,counter_pdf,counter_docx,counter_xlsx,counter_pptx
    
   
    if file.endswith(".txt"):
        file_path=os.path.join(root, file)
        data=''
        
        with open(file_path, 'r') as f:
            data+=f.read()
            result=re.findall(text,data)
            if  result:
                final.append(result[0]+"in file"+file)
        counter_txt+=1

                    
                
    if file.endswith(".pdf"):
        file_path=os.path.join(root, file)
        doc = fitz.open(file_path)
        data=''
        for page in doc:
            data+=page.get_text()
        result=re.findall("data",data)
        if result:
            final.append(result[0]+" in file "+file)
        counter_pdf+=1
            

    if file.endswith(".xlsx"):
        file_path=os.path.join(root, file)
        data=''
        with open(file_path, 'rb') as f:
            wb = load_workbook(f, read_only=True, data_only=True)
            ws = wb.active
            
            values=[str(cell.value) for row in ws for cell in row if cell.value is not None]
            data=' '.join(values)            
            result=re.findall(text,data)
            if result:
                final.append(result[0]+" in file "+file)
        counter_xlsx+=1
            

    if file.endswith(".pptx"):
        file_path=os.path.join(root, file)
        data=''
        prs = Presentation(file_path)
        values=[shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
        data=' '.join(values)
        result=re.findall(text,data)
        if result:
            final.append(result[0]+" in file "+file)
        counter_pptx+=1

    if file.endswith(".docx"):
        file_path=os.path.join(root, file)
        data=''
        doc = Document(file_path)
        data=''
        for para in doc.paragraphs:
            data += '\n'+ para.text
        result=re.findall(text,data)
        if result:
            final.append(result[0]+" in file "+file)
        counter_docx+=1
    counter+=1
    

#Going through each file in the directory to unpack
[extract(root,file) for root, dirs, files in os.walk(directory) for file in files if file.endswith((".zip",".rar")) ]#map(str.upper,os.walk(directory))


[check(root,file)  for root, dirs, files in os.walk(directory) for file in files if file.endswith((".txt",".pdf",".xlsx",".pptx",".docx"))]
print("\n")
[print(r) for r in final]

executionTime = (time.time() - startTime)
print('Execution time in seconds: ' + str(executionTime))

if out_zip !='' and os.path.exists(out_zip) ==True:
    ans = input('Do you want to remove all the files that were extracted from compressed directories?   y/n?')
    if ans.lower() =='y':
        shutil.rmtree(out_zip)
        


if args.save_log is not None:
    log_dir_path=os.path.join(directory, "search_log")
    if not os.path.exists(log_dir_path):
        os.makedirs(log_dir_path)

    now = datetime.datetime.now()
    log_path=os.path.join(log_dir_path, now.strftime('%Y_%m_%d__%H_%M_%S')+" search file log for "+text+".txt")
    with open(log_path,"w") as logtxt:
        log_data=[logsize(root,file) for root, dirs, files in os.walk(directory) for file in files if file.endswith((".txt",".pdf",".xlsx",".pptx",".docx")) ]
        logtxt.write(f"Total text files {counter_txt}\nTotal pdf files {counter_pdf}\nTotal docx files {counter_docx}\nTotal xlsx files {counter_xlsx}\nTotal pptx files {counter_pptx}\n")
        x=[len(dirs) for root, dirs, files in os.walk(directory)]
        logtxt.write(f"Total directories {x[0]}\n")
        logtxt.write(f"Total files that are not searched {counter_otherfiles}\n")
        for data in log_data:
            logtxt.write(data+"\n")
    

