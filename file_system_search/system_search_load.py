#This script goes through every sub folder from an identified folder
#It searches for a string within popular document types: office, pdf, archives

import os,re
import time
import datetime
import argparse
import shutil

try:
    import fitz
    import pandas as pd
    import numpy as np
    from pptx import Presentation
    from docx import Document
    from zipfile import ZipFile
    from rarfile import RarFile
    import shutil
except ModuleNotFoundError:
    install = input("Requirments not found on computer, would you like to install them automatically? Y/N\n>>> ")
    if install.lower() == 'y':
        os.system("pip install pymupdf")
        os.system("pip install pandas")
        os.system("pip install numpy")
        os.system("pip install pptx")
        os.system("pip install docx")
        os.system("pip install zipfile39")
        os.system("pip install rarfile")





startTime = time.time()
counter,totalsize,counter_txt,counter_pdf,counter_docx,counter_xlsx,counter_pptx,counter_totalfiles,counter_otherfiles,count_not,searchsize=0,0,0,0,0,0,0,0,0,0,0
final=list()
out_zip=''
not_searched=[]

#Setting up arguments for command line
parser = argparse.ArgumentParser(prog="system_search_load.py")
parser.add_argument('-path', '--source_path', help="Takes the directory path in quotes as input")
parser.add_argument('-text', '--text_pattern', help="Takes the text pattern as input")
parser.add_argument('save', nargs='?', help="Saves a log file in a folder")

args = parser.parse_args()

#Assigning variables and adding regular expression for the text
directory = args.source_path
text = args.text_pattern+".+" #Has the text with one or more characters and ends with }


def logsize(root,file):
    global totalsize
    size=os.path.getsize(os.path.join(root, file))
    totalsize+=size
    return "{f} | {size} | Bytes".format(f=file,size=size)

def getsize(root,file):
    size=os.path.getsize(os.path.join(root, file))
    return "{size} | Bytes".format(size=size)

def other(file):
    global counter_otherfiles
    counter_otherfiles+=1

def total_files(file):
    global counter_totalfiles
    counter_totalfiles+=1


def progress(percent=0,total=0, width=50):
    left = width * percent // total
    right = width - left
    
    tags = "#" * left
    spaces = " " * right
    percents = f"{(percent/total)*100:.1f}% done | {percent} files done"
    
    print("\r[", tags, spaces, "]", percents, sep="", end="", flush=True)
# Example run

def extract(root,file):
    global out_zip
    try:
        out_zip=os.path.join(root, "extracted")
        file_path=os.path.join(root, file)
        if file.endswith(".zip"):
            with ZipFile(file_path,"r") as z:
                z.extractall(out_zip)
        if file.endswith(".rar"):
            with RarFile(file_path,"r") as r:
                r.extractall(out_zip)
    except Exception as e:
        print(e)
    
def check(root,file,total_files):
    global counter,counter_txt,counter_pdf,counter_docx,counter_xlsx,counter_pptx,count_not,not_searched,searchsize
    
    try:
    
        if file.endswith(".txt"):
            file_path=os.path.join(root, file)
            data=''
            
            with open(file_path, 'r', encoding="utf8") as f:
                data+=f.read()
                result=re.findall(text,data)
                if  result:
                    final.append([result[0],file_path,getsize(root,file)])
            
            counter_txt+=1
              
        if file.endswith(".pdf"):
            file_path=os.path.join(root, file)
            doc = fitz.open(file_path)
            data=''
            for page in doc:
                data+=page.get_text()
            result=re.findall(text,data)
            if result:
                final.append([result[0],file_path,getsize(root,file)])
            counter_pdf+=1
    
        if file.endswith(".xlsx"):
            file_path=os.path.join(root, file)
            data=''
            df = pd.read_excel(file_path, engine='openpyxl')            
            mask = np.column_stack([df[col].astype(str).str.contains(text, na=False) for col in df])
            indices = np.argwhere(mask)
            if len(indices)>0:
            
                final.append("Found in file "+file)
            counter_xlsx+=1
    
        if file.endswith(".pptx"):
            file_path=os.path.join(root, file)
            data=''
            prs = Presentation(file_path)
            values=[shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
            data=' '.join(values)
            result=re.findall(text,data)
            if result:
                final.append([result[0],file_path,getsize(root,file)])
            counter_pptx+=1
    
        if file.endswith(".docx"):
            file_path=os.path.join(root, file)
            data=''
            doc = Document(file_path)
            data=''
            for para in doc.paragraphs:
                data += '\n'+ para.text
            result=re.findall(text,data)
            if result:
                final.append([result[0],file_path,getsize(root,file)])
            counter_docx+=1
        searchsize+=os.path.getsize(os.path.join(root, file))
        counter+=1
    except Exception as e:
        osize=os.path.getsize(os.path.join(root, file))
        count_not+=1
        not_searched.append(str(file)+" | "+str(osize)+" | Bytes")
        not_searched.append(str(e))
        

    
    progress(counter,total_files)
    

#Going through each file in the directory to unpack
[extract(root,file) for root, dirs, files in os.walk(directory) for file in files if file.endswith((".zip",".rar")) ]#map(str.upper,os.walk(directory))

[total_files(file) if file.endswith((".txt",".pdf",".xlsx",".pptx",".docx")) and not file.startswith("~$") else other(files) for root, dirs, files in os.walk(directory) for file in files ]


[check(root,file,counter_totalfiles)  for root, dirs, files in os.walk(directory) for file in files if file.endswith((".txt",".pdf",".xlsx",".pptx",".docx")) and not file.startswith("~$")]
print("\n")
[print(value) for values in final for value in values]
                    

executionTime = (time.time() - startTime)
print('Execution time in seconds: ' + str(executionTime))

if out_zip !='' and os.path.exists(out_zip) ==True:
    ans = input('Do you want to remove all the files that were extracted from compressed directories?   y/n?')
    if ans.lower() =='y':
        shutil.rmtree(out_zip)
        
if args.save is not None:
    log_dir_path=os.path.join(directory, "search_log")
    if not os.path.exists(log_dir_path):
        os.makedirs(log_dir_path)

    now = datetime.datetime.now()
    log_path=os.path.join(log_dir_path, now.strftime('%Y_%m_%d__%H_%M_%S')+" search file log for "+text+".txt")
    with open(log_path,"w", encoding="utf-8") as logtxt:
        total=0
        for root, dirs, files in os.walk(directory):
            for file in files:
                total+=os.path.getsize(os.path.join(root, file))    

        log_data=[logsize(root,file) for root, dirs, files in os.walk(directory) for file in files if file.endswith((".txt",".pdf",".xlsx",".pptx",".docx")) and not file.startswith("~$") ]
        logtxt.write("Numeric logs: \n\n")
        logtxt.write(f"Total text files: {counter_txt}\nTotal pdf files: {counter_pdf}\nTotal docx files: {counter_docx}\nTotal xlsx files: {counter_xlsx}\nTotal pptx files: {counter_pptx}\n")
        logtxt.write("Total files size :{size} MB\n".format(size=round(total/1024**2,3)))
        x=[len(dirs) for root, dirs, files in os.walk(directory)]
        logtxt.write(f"Total directories: {x[0]}\n")
        logtxt.write(f"Total files that are not searched: {counter_otherfiles}\n")
        logtxt.write("Total searched files size: {ssize} MB\n".format(ssize=round(totalsize/1024**2,3)))
        logtxt.write('Execution time in seconds: ' + str(executionTime)+" seconds\n\n")
        logtxt.write("RESULT of searched text\n\n")
        if not final:
            logtxt.write('Did not find any text in files \n\n')
        else:
            for values in final:
                for value in values:
                    logtxt.write(value+"\n")
        logtxt.write("\n\n")
        logtxt.write('ENCOUNTERERED files and thier reason: \n\n')
        for val in not_searched:
            logtxt.write(val+"\n")
        logtxt.write("\n")
        logtxt.write('Listed searched files and thier size: \n\n')
        for data in log_data:
            logtxt.write(data+"\n")
    

