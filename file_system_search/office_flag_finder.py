import sys,os,re
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document

file = sys.argv[1]
text= sys.argv[2]+".+}"

if file.endswith(".xlsx"):
    data=''
    wb = load_workbook(file)
    ws = wb.active
    for row in ws:
        for cell in row:
            if cell.value is not None:
                data+=str(cell.value)
    result=re.findall(text,data)
    print(result[0])
                   

if file.endswith(".pptx"):
    
    data=''
    with open(file, 'rb') as f:
        prs = Presentation(f)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    data=shape.text
        result=re.findall(text,data)
        print(result[0])
        

if file.endswith(".docx"):
    data=''
    with open(file, 'rb') as f:
        doc = Document(file)
        data=''
        for para in doc.paragraphs:
            data += '\n'+ para.text
        result=re.findall(text,data)
        print(result[0])
        